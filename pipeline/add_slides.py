"""Append two slides to SEEM_06.02.2025.pptx in the deck's own style."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

SCRATCH = '/private/tmp/claude-501/-Users-tom-Repos-SEEM/a72dfa34-e648-4027-ba2b-3b7da1981ec8/scratchpad'
PATH = '/Users/tom/Downloads/SEEM_06.02.2025.pptx'
DARK = RGBColor(0x26, 0x26, 0x26)
GREY = RGBColor(0x40, 0x40, 0x40)
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def box(slide, x, y, w, h):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    b.text_frame.word_wrap = True
    return b


def style_para(para, bullet):
    pPr = para._p.get_or_add_pPr()
    for tag in ('spcAft', 'buSzPct', 'buChar', 'buNone'):
        for el in pPr.findall(f'{{{NS}}}{tag}'):
            pPr.remove(el)
    if bullet:
        pPr.set('marL', '342900')
        pPr.set('indent', '-342900')
    spc = etree.SubElement(pPr, f'{{{NS}}}spcAft')
    etree.SubElement(spc, f'{{{NS}}}spcPts').set('val', '800')
    if bullet:
        etree.SubElement(pPr, f'{{{NS}}}buSzPct').set('val', '100000')
        etree.SubElement(pPr, f'{{{NS}}}buChar').set('char', '•')
    else:
        etree.SubElement(pPr, f'{{{NS}}}buNone')


def put(tf, lines, size=16, bold=False, italic=False, color=GREY,
        font='Calibri', bullet=False):
    first = True
    for line in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = para.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        style_para(para, bullet)


prs = Presentation(PATH)
layout = prs.slides[1].slide_layout

# ---------------- slide 3: symbolic regression ----------------
s = prs.slides.add_slide(layout)
put(box(s, 0.70, 0.55, 12.0, 0.80).text_frame,
    ['Symbolic Regression (searching the space of formulas)'],
    size=30, bold=True, color=DARK)
put(box(s, 0.70, 1.50, 12.0, 0.40).text_frame,
    ['A regression method where the unknown is the formula itself — its '
     'structure and its constants are found together:'],
    size=16, color=GREY)
put(box(s, 0.90, 2.00, 11.6, 0.50).text_frame,
    ['minimize  (fit error)²  over all formulas f : complexity(f) ≤ C'],
    size=20, italic=False, color=DARK, font='Consolas')
put(box(s, 0.90, 2.90, 11.6, 1.85).text_frame,
    ['Formulas are trees built from x, fitted constants, and the operators '
     'exp, log, 1/x, x², √x, +, −, ×, ÷.',
     'Instead of sampling this space (genetic programming), our engine '
     'enumerates every distinct formula up to the budget — 624,936 trees '
     'per fit at C = 8 — so "no better formula exists in this space" is '
     'a checkable claim, not luck.',
     'Constants are optimized inside the search (Levenberg–Marquardt, '
     '4 restarts, outer scale/offset solved exactly); unphysical candidates '
     '(non-monotonic, divergent) are filtered out; the output is the exact '
     'accuracy-vs-complexity Pareto front.'],
    size=16, color=GREY, bullet=True)
put(box(s, 0.70, 4.85, 12.0, 0.40).text_frame,
    ['The role of the complexity budget (C):'], size=20, bold=True, color=DARK)
put(box(s, 0.90, 5.35, 11.6, 1.30).text_frame,
    ['Small C  →  simple, interpretable formulas that may miss real '
     'structure.',
     'Large C  →  flexible formulas that can memorize noise — so we '
     'validate on synthetic data, cross-validate by whole trajectory, and '
     'spot-check one budget higher (4,505,024 trees at C = 9: nothing new).'],
    size=16, color=GREY, bullet=True)

# ---------------- slide 4: distribution-fit results ----------------
s = prs.slides.add_slide(layout)
put(box(s, 0.70, 0.55, 12.0, 0.70).text_frame,
    ['Result: the yield-event size distribution'],
    size=30, bold=True, color=DARK)
put(box(s, 0.70, 1.30, 12.0, 0.40).text_frame,
    ['Symbolic regression on 187,468 stress-drop events, run cell-by-cell in '
     'the state plane (u, τ); winners adjudicated by maximum likelihood.'],
    size=16, italic=True, color=GREY)
put(box(s, 0.90, 1.95, 11.8, 1.70).text_frame,
    ['The textbook truncated power law s⁻ᵏ exp(−s/s*) '
     'never appears on any Pareto front. The discovered law is '
     'P(s) ∝ (s_c − s)ᵐ / (s + ε)ᵏ  with k = 0.88, '
     'm = 2.3 — a power law that ends at a hard ceiling s_c.',
     'One shape, one state-dependent number: all 12 test cells collapse onto '
     'a single curve, and the ceiling field s_c(u, τ) — spanning '
     '146× — carries all the state (and history) dependence.',
     'The verdict is falsifiable: on synthetic exponential-tail data the same '
     'pipeline returns the power law (5/5 seeds); on synthetic ceiling data '
     'it recovers the ceiling (5/5). The measured data lands on the ceiling '
     'side, ΔAIC ≈ 200–345 per cell.'],
    size=16, color=GREY, bullet=True)
s.shapes.add_picture(SCRATCH + '/panel_a.png', Inches(1.30), Inches(3.80),
                     height=Inches(3.45))
s.shapes.add_picture(SCRATCH + '/panel_b.png', Inches(7.35), Inches(3.80),
                     height=Inches(3.45))

prs.save(PATH)
print('saved; slides now:', len(prs.slides))
