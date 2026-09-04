import os as _os
"""SEEM project deck v3: the distribution-first pass and Section 9, in the plain research format of v2."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1E, 0x27, 0x61)
INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x6E, 0x7B, 0x85)
RED = RGBColor(0x99, 0x00, 0x11)
BOX = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
P = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..', 'figures') + '/'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, x, y, w, h, wrap=True, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def run(p, text, size=15, bold=False, color=INK, font='Calibri'):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def set_bullet(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '228600')
    pPr.set('indent', '-228600')
    bu = pPr.makeelement(qn('a:buChar'), {'char': '•'})
    pPr.append(bu)


def bullets(tf, items, size=15):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, hdr = it
        else:
            txt, hdr = it, False
        p = para(tf, first=(i == 0))
        if hdr:
            run(p, txt, size=size + 1, bold=True, color=NAVY)
            p.space_after = Pt(6)
        else:
            run(p, txt, size=size, color=INK)
            set_bullet(p)
            p.space_after = Pt(8)


def title(s, t, color=NAVY):
    tf = textbox(s, 0.6, 0.32, 12.13, 0.75)
    run(para(tf, True), t, size=30, bold=True, color=color)


def takeaway(s, t):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.62))
    sh.adjustments[0] = 0.18
    sh.fill.solid()
    sh.fill.fore_color.rgb = BOX
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run(p, 'Takeaway: ', size=13.5, bold=True, color=NAVY)
    run(p, t, size=13.5, color=INK)


def img(s, path, x, y, w, h):
    s.shapes.add_picture(P + path, Inches(x), Inches(y), Inches(w), Inches(h))


def caption(s, t, y=5.75, color=INK, size=13.5, align=None):
    tf = textbox(s, 0.7, y, 12.0, 0.7)
    p = para(tf, True)
    run(p, t, size=size, color=color)
    if align:
        p.alignment = align

def img_fit(s, path, x, y, w, h):
    """place the image inside the (x,y,w,h) box, preserving aspect, centred"""
    from PIL import Image
    iw, ih = Image.open(P + path).size
    sc = min(w / iw, h / ih)
    ww, hh = iw * sc, ih * sc
    s.shapes.add_picture(P + path, Inches(x + (w - ww) / 2), Inches(y + (h - hh) / 2), Inches(ww), Inches(hh))


def fig_left(s, path, items, size=14, w=6.9, h=4.95, y=1.35, cap=None):
    img_fit(s, path, 0.55, y, w, h)
    tf = textbox(s, 0.55 + w + 0.35, y + 0.15, 12.78 - w - 0.35, h - 0.2)
    bullets(tf, items, size=size)
    if cap:
        caption(s, cap, y=6.05, size=12.5, color=MUTED)


def fig_wide(s, path, items, size=13.5, h=3.45, y=1.3):
    img_fit(s, path, 0.55, y, 12.2, h)
    tf = textbox(s, 0.7, y + h + 0.12, 12.0, 6.5 - (y + h + 0.12))
    bullets(tf, items, size=size)


def two_figs(s, a, b, cap=None, y=1.3, h=4.3):
    img_fit(s, a, 0.55, y, 6.0, h)
    img_fit(s, b, 6.78, y, 6.0, h)
    if cap:
        caption(s, cap, y=y + h + 0.12, size=13)


# ------------------------------------------------------------------ 1 title
s = slide()
tf = textbox(s, 0.9, 1.9, 11.5, 1.6)
run(para(tf, True), 'A fitted law for yield events, and everything rebuilt from it', size=38, bold=True, color=NAVY)
tf = textbox(s, 0.9, 3.55, 11.5, 0.9)
run(para(tf, True), 'Symbolic-regression discovery of the event-size distribution in AQS shear, '
                    'the forward model built on it, and where the model fails', size=19, color=INK)
tf = textbox(s, 0.9, 4.6, 11.5, 0.45)
run(para(tf, True), 'SEEM project status, September 4, 2026  (third pass; supersedes the August 6 deck)', size=15, color=MUTED)
tf = textbox(s, 0.9, 6.5, 11.5, 0.4)
run(para(tf, True), 'AQS molecular-dynamics shear · 850 trajectories · nine preparations · '
                    '187,468 stress drops · every number regenerated by pipeline/run_all.sh', size=13, color=MUTED)

# ---------------------------------------------------------------- 2 dataset
s = slide()
title(s, 'The dataset, and the object of study')
for i, (num, lab) in enumerate([('850', 'trajectories'), ('9', 'cooling rates (256×)'),
                                ('187,468', 'stress-drop events'), ('105,188', 'aging events')]):
    x = 0.7 + i * 3.05
    tf = textbox(s, x, 1.25, 2.9, 0.85); p = para(tf, True)
    run(p, num, size=36, bold=True, color=NAVY); p.alignment = PP_ALIGN.CENTER
    tf = textbox(s, x, 2.08, 2.9, 0.35); p = para(tf, True)
    run(p, lab, size=13, color=MUTED); p.alignment = PP_ALIGN.CENTER
tf = textbox(s, 0.7, 2.85, 5.9, 3.5)
bullets(tf, [('Protocol', True),
             'Athermal quasistatic shear of a model glass; Δγ = 10⁻⁵, 49,999 steps per run, γ → 0.5',
             'Cooling rates 2×10¹⁰ → 5.12×10¹² K/s in factor-of-2 steps, ~100 samples each',
             'Per step: strain, stress, potential energy, cooling rate. Nothing else is recorded'])
tf = textbox(s, 6.95, 2.85, 5.8, 3.5)
bullets(tf, [('State and events', True),
             'u = pe − u₀ (energy above the reference), τ = stress/10⁴',
             'Stress drop: a step with dτ < 0, size s = −dτ;  aging: du < 0 without a drop',
             '22 × 22 grid on (u,τ): inner 20 × 20 on the 0.2 to 99.8% quantiles plus a catch-all ring'])
takeaway(s, 'The target this time is the probability distribution of yield events, P(s | u,τ), with its form '
            'discovered rather than assumed, and everything downstream is rebuilt from it.')

# -------------------------------------------------------------- 3 framework
s = slide()
title(s, 'Framework: the step process as drift + jumps')
tf = textbox(s, 0.7, 1.35, 5.9, 4.9)
bullets(tf, [('Three fields measured directly, per cell', True),
             'Elastic drift 2μ(u,τ), and the energy drift on quiet steps',
             'Event hazard p(u,τ): fraction of steps that drop (median 0.32% per step)',
             'Energy released per event: du = a + b·s + noise, per cell; an aging channel alongside'])
tf = textbox(s, 6.95, 1.35, 5.8, 4.9)
bullets(tf, [('One law to discover', True),
             'The size distribution P(s | u,τ): four decades of s, 187k events',
             'First pass assumed a truncated power law (TPL); this pass asks the data for the form',
             'Everything downstream (mean sizes, the plastic-rate field q, full synthetic stress-strain curves) '
             'is then derived from the fitted law and compared with its measured counterpart'])
takeaway(s, 'Where derived and measured agree, the law carries the physics; where they disagree, we say so.')

# -------------------------------------------------------------- 4 SR engine
s = slide()
title(s, 'Symbolic regression, built for scrutiny')
tf = textbox(s, 0.7, 1.35, 5.9, 4.9)
bullets(tf, [('Engine  (library/symreg.py)', True),
             'Exhaustive enumeration of expression trees: unary exp, log, 1/x, x², √x, −x; binary + − × ÷',
             'Up to two constants fitted inside the nonlinearity (Levenberg-Marquardt, 4 restarts); outer affine profiled out',
             'Numerical fingerprint dedup; hard admissibility (monotone decreasing); exact Pareto fronts',
             'Complexity ≤ 8: 624,936 trees → ~14,000 distinct functions → ~10,300 admissible per cell'])
tf = textbox(s, 6.95, 1.35, 5.8, 4.9)
bullets(tf, [('Protocol on the distributions', True),
             'Target: binned ln ρ(s) per cell, 19 log bins, Poisson errors',
             '12 cells (700 to 7,600 events) spanning the plane, plus unweighted-bin and wider-window variants',
             'Audit: the largest cell pushed to complexity 9 (4,505,024 trees, 48,725 admissible); no new form appears',
             'The search proposes forms; a per-event likelihood chooses among them (next slides)'])
takeaway(s, 'Search the whole grammar, admit only monotone-decreasing densities, and let a per-event likelihood '
            'adjudicate among the Pareto knees.')

# ---------------------------------------------------------- 5 law in one cell
s = slide()
title(s, 'The discovered law, in one cell')
fig_left(s, 'fig01_size_law.png', [
    'The SR knee in the largest cells is log(1/(s+ε) − C), which rearranges to a bounded-support form:',
    'P(s | u,τ) ∝ (s_c − s)ᵐ / (s + ε)ᵏ  on 0 < s < s_c',
    'k = 0.88: the small-event power law;  ε = 2.8×10⁻⁵: rounds it off below the resolution',
    's_c(u,τ): a hard ceiling, the largest event the state can produce; the density vanishes like (s_c − s)ᵐ, m ≈ 2.3',
    'The TPL’s exponential tail merely suppresses large events; the ceiling forbids them'], size=13.5)
takeaway(s, 'Four decades of event size in one cell follow one four-parameter law whose only state-dependent '
            'parameter is the ceiling s_c.')

# --------------------------------------------------------- 6 collapse + pareto
s = slide()
title(s, 'One shape everywhere; the fronts that found it')
two_figs(s, 'fig02_collapse.png', 'fig03_pareto.png',
         'Left: twelve cells spanning 100× in scale collapse onto one curve in s/s_c with global (k, m, ε).  '
         'Right: per-cell Pareto fronts (heights set by each cell’s own noise floor); orange, the complexity-9 audit run.')
takeaway(s, 'The state enters only through the ceiling: the shape is universal across the (u,τ) plane.')

# ------------------------------------------------------- 7 likelihood disposes
s = slide()
title(s, 'The likelihood chooses among the candidate forms')
fig_left(s, 'fig04_falsifiability.png', [
    'Normalized maximum likelihood per cell, seven families, AIC and trajectory-blocked cross-validation',
    'The ceiling law wins 8 of 8 large cells: ΔAIC 67 to 120 over the TPL, 300 to 754 over the lognormal',
    'A binned-least-squares favourite (√log form) was an artifact of Poisson weighting: ΔAIC ≈ +20,000 under likelihood',
    'Synthetic control, both directions: TPL-generated catalogs return a TPL verdict 5/5; '
    'ceiling-generated return a ceiling verdict 5/5 with ε and s_c recovered to a few %',
    'Sparse cells (< 1,000 events) cannot tell a ceiling from an exponential tail'], size=13)
takeaway(s, 'Per-event likelihood with falsifiability controls decides for the bounded-support law; the binned fit only proposed it.')

# ------------------------------------------------------------- 8 ceiling field
s = slide()
title(s, 'The ceiling field s_c(u,τ)')
two_figs(s, 'fig05_ceiling_map.png', 'fig07_ceiling_memory.png',
         'Left: 277 cells, 146× in range (0.02 cold → 3.0 in the flowing state), bootstrap error 6.5%; '
         'rank-1 factorization Λ(u)·f(τ) carries 95.5% of the log-variance.  '
         'Right: refit by preparation at fixed cell; slow-cooled glass has a 1.71× higher ceiling than fast-cooled.')
takeaway(s, 'All state dependence, and all preparation memory, sit in one parameter: the ceiling.')

# ------------------------------------------------------------- 9 q rebuilt
s = slide()
title(s, 'Rebuilding the plastic-rate field from the law')
fig_left(s, 'fig09_q_agreement.png', [
    'q = 1 − ⟨dτ⟩ / (2μ Δγ), measured per cell in the first pass',
    'From the law: hazard × mean event size of the fitted distribution',
    'Agreement: median +1.8%, RMS 17% against a 17.5% split-half noise floor; the test is saturated',
    'The Λ(u)·f(τ) factors of the rebuilt field are the measured ones: '
    'f(τ) ∝ τ·log(1/τ − 0.44) vs 0.43; the Λ pole at u = 0.0287 vs 0.0280',
    'Old κ ≈ 0.95 recovered as the mid-range slope; the x_min drift problem is gone with ε'], size=13.5)
takeaway(s, 'The fitted law reproduces the first pass’s central object to within its noise floor.')

# ------------------------------------------------------------ 10 forward sim
s = slide()
title(s, 'Simulating all 850 runs')
two_figs(s, 'fig12_sim_stress.png', 'fig14_sim_peaks.png',
         'Three arms differ only in the size law: ceiling law, TPL, empirical resampling.  Stress-strain curves, '
         'yield overshoot across 256× in cooling rate, and flow stress (±5%) all emerge; '
         'yield-peak RMS 2.5% (ceiling), 1.4% (TPL), 0.9% (empirical).')
takeaway(s, 'The model predicts the yield strength of a glass from its preparation to a few percent, '
            'and macroscopic curves cannot tell size laws apart. The per-event likelihood can.')

# ----------------------------------------------------------- 11 failure mode
s = slide()
title(s, 'The failure mode', color=RED)
fig_left(s, 'fig13_sim_energy.png', [
    'The energy coordinate: good early, too convergent late',
    'Preparations keep a 24% spread in mean u at γ = 0.5; the simulation keeps 4%',
    'The fields are preparation-blind by construction (one map for all histories), '
    'so this is the missing state variable of the first pass, seen macroscopically',
    'Same message as the 1.71× ceiling memory: (u,τ) is not a sufficient state'], size=14)
takeaway(s, 'A five-times-too-small preparation spread: the one place the two-variable model fails, and the question '
            'the rest of the deck answers.')

# --------------------------------------------------------- 12 where memory lives
s = slide()
title(s, 'Where the preparation memory lives')
fig_left(s, 'fig16_memory_channels.png', [
    'Every channel re-measured per cooling rate at fixed cell; slope β on ln(rate) with cell fixed effects; '
    'a parity-split control gives the null',
    'Event hazard: preparation-blind (β = −0.012 ± 0.002; ratio 1.005, control 1.02). Elastic drift: identical.',
    'Mean event size: 1.70× (β = −0.117 ± 0.005), the ceiling memory measured another way; '
    'it lives in the tail (geometric mean moves only 1.20×)',
    'A slow-cooled glass is triggered as often; its avalanches run further',
    '73% of the per-event energy-release memory is size; 27% a secondary memory in the coupling'], size=12.5)
takeaway(s, 'Memory acts through event size, not event rate. That cuts against the simplest one-variable STZ reading, '
            'where preparation enters through the STZ density.')

# ------------------------------------------------------------- 13 repair
s = slide()
title(s, 'A one-parameter memory in the ceiling repairs the failure')
fig_wide(s, 'fig17_memory_repair.png', [
    'Single change: s_c(u,τ; rate) = s_c(u,τ)·(rate/rate_mid)^β with the one global β = −0.117; hazard, drift, aging untouched',
    'Spread of mean u at γ = 0.5: 4.1% → 26.5% against 24% in the data; the curve tracks the data over the whole strain window; '
    'yield peaks unchanged (2.5% → 2.8%)',
    'Refinements overshoot (β by u-band: 37%; per-preparation coupling: 35%); the simple statement is the defensible one'], size=12.5, h=3.9)
takeaway(s, 'To the resolution of this dataset, the missing state variable is a single multiplicative factor on the ceiling.')

# ------------------------------------------------------------- 14 relaxes
s = slide()
title(s, 'The memory relaxes with strain')
fig_wide(s, 'fig19_memory_relaxes.png', [
    'Nine-rate regression with a fixed effect per (cell, strain window): β = −0.21 → −0.11 → −0.05 → −0.04 from γ ≈ 0.15 to 0.45, '
    'and it falls within fixed u bands, so this is relaxation rather than state dependence; γ_r ≈ 0.18',
    'Late strain: size memory 1.75× → 1.16×, hazard blind throughout; the energy-step memory only halves, '
    'a third of it in the coupling at fixed size',
    'Whether that coupling memory also relaxes needs runs past γ = 0.5'], size=12.5, h=3.9)
takeaway(s, 'The third variable is not a frozen label: plastic strain erases it on a scale of a few tenths, '
            'now measured, not assumed.')

# ------------------------------------------------------------- 15 per-sample
s = slide()
title(s, 'The slow variable, one trajectory at a time')
fig_left(s, 'fig22_persample.png', [
    'Parameter-free prediction: apply the between-preparation relation within a preparation, so '
    'each trajectory gets exp(κ·δu₀) with κ = β·d ln(rate)/du₀ = −43',
    'Predicted persistence of a sample’s initial energy deviation: 0.44 / 0.27 / 0.22 / 0.17 at γ = 0.1 / 0.2 / 0.3 / 0.5',
    'Data: 0.63 / 0.27 / 0.24 / 0.18.  Base model: 0.41 / 0.08 / 0.07 / 0.08.  Tripling κ overshoots.',
    'Same variable, same strength, acting on individual samples'], size=13.5)
takeaway(s, 'A sample remembers its initial energy exactly as long as the preparation-level ceiling memory says it should.')

# ---------------------------------------------------------- 16 ceiling robustness
s = slide()
title(s, 'Is the ceiling just the largest event?')
fig_left(s, 'fig18_ceiling_pinning.png', [
    'Refit with the top 1, 5, or 1% of events removed: s_c tracks the new maximum almost one-for-one (red)',
    'So does a true ceiling with m ≈ 2 at n ≈ 5,000 (green synthetics: c/s_max 1.02 to 1.12)',
    'A true TPL under the same fit runs off to 3 to 350× the maximum and collapses on refit (blue)',
    'Merging consecutive drop steps into single avalanches (187k → 172k): the ceiling law wins by more '
    '(ΔAIC 117 to 213), s_c moves < 2%, k shifts +0.05',
    's_c in a cell is set by its ~10 largest events: the 6.5% bootstrap error is a floor'], size=13)
takeaway(s, 'The data sit on the true-ceiling line, not the exponential cloud, and the law survives the event definition.')

# ------------------------------------------------------------- 17 aftershocks
s = slide()
title(s, 'Where the Markov assumption fails: aftershocks')
fig_wide(s, 'fig20_clustering.png', [
    'At the same cell, the hazard on the step after an event is 20× the quiet baseline, decaying over ~5 steps; '
    '14% of events sit in consecutive runs (Markov expectation 0.4%); the next event is 1.2× larger',
    'Event counts per window are over-dispersed (Fano 1.3 to 1.6 vs the model’s 1.0) while the stress released per window is less variable (CV 0.10 vs 0.13)',
    'Referenced to the last large event, the large-event hazard is 1.5× higher after 0.1 to 0.3 of reloading; large events are more regular (Fano 0.51 vs 0.71)'], size=12.5, h=3.9)
takeaway(s, 'No function of (u,τ) produces this; an aftershock term fixes the yield peaks (RMS 2.5% → 1.6%) but not the variance below.')

# ------------------------------------------------------------- 18 variance gap
s = slide()
title(s, 'The variance gap: ten exclusions, one open problem')
img_fit(s, 'fig21_renewal.png', 0.55, 1.3, 7.4, 2.85)
tf = textbox(s, 8.15, 1.3, 4.65, 2.9)
bullets(tf, [('Run-to-run scatter within a preparation: model 25 to 30% too high', True),
             'Not the size law, aftershock hazard, grid (42×42), renewal hazard, two-class reload hazard',
             'Not size-reload coupling (+0.02), size-size coupling (−0.01), relaxation rate (identical), '
             'the per-sample slow variable, or the u-channel noise terms'], size=12)
tf = textbox(s, 0.7, 4.3, 12.0, 2.2)
bullets(tf, ['Each candidate was run forward and scored on the same three metrics (count Fano factor, CV of stress released per window, τ scatter)',
             'What survives: a driving-noise amplitude 17% too large with the same restoring force, at fixed (u,τ), '
             'that no sequence correlation the data can exhibit accounts for',
             'Closing it needs a per-step observable that is not (u,τ), the same MD output the third variable asks for'], size=12.5)
takeaway(s, 'The remaining discrepancy is bounded and named; the dataset has exhausted what it can test at fixed (u,τ).')

# ------------------------------------------------------------------ 19 verdict
s = slide()
title(s, 'What holds, what doesn’t')
def col(x, head, colr, items):
    tf = textbox(s, x, 1.35, 3.9, 0.45)
    run(para(tf, True), head, size=17, bold=True, color=colr)
    tf = textbox(s, x, 1.9, 3.9, 4.5)
    bullets(tf, items, size=12.5)
col(0.7, 'SUPPORTED', NAVY, [
    'One universal size shape with a single state-dependent scale (12 cells, 100× in scale)',
    'A hard ceiling, sharper than exponential (8/8 by AIC; synthetic controls both ways; robust to merged avalanches)',
    'Small-size exponent k ≈ 0.9; the ceiling field factorizes (95.5%)',
    'Stress-strain curves and yield peaks across 256× in cooling rate at the few-% level',
    'Memory in event size only; relaxes with γ_r ≈ 0.2; parameter-free per-sample prediction'])
col(4.85, 'NOT ESTABLISHED', MUTED, [
    'Ceiling exponent m ≈ 2: order of magnitude only',
    'Specific SR forms of the ceiling factors: best-of-front, not unique',
    'Ceiling vs exponential in cells with < 1,000 events',
    'Convergence of preparations at large strain: the model predicts it, now with a measured rate, but there is no data past γ = 0.5',
    'The system size N (unrecorded): the finite-size reading of the ceiling is untested'])
col(9.0, 'CONTRADICTED', RED, [
    '(u,τ) as a sufficient state: the 1.71× ceiling memory and the 5×-underpredicted u-spread, independently',
    'The step process as Markov in (u,τ) at the 10⁻⁵ strain scale: aftershocks (×20), reload-dependent large events',
    'Any of ten candidate explanations for the run-to-run variance excess',
    'Preparation acting through the event rate (the simplest STZ-density reading)'])
takeaway(s, 'The law and the model built on it hold at the few-percent level; their failures are located, quantified, and '
            'point at one slow variable and one non-Markov structure.')

# --------------------------------------------------------------- 20 next steps
s = slide()
title(s, 'Recommended next steps')
steps = [
    ('A per-step structural observable from the MD', 'written alongside pe (local order, soft modes): the one test is whether it absorbs the '
     'ceiling’s β = −0.117 while leaving the hazard’s β at zero; it is also the last lever on the variance gap'),
    ('Runs past γ = 0.5', 'one preparation to γ ≈ 2 settles whether the coupling memory relaxes like the ceiling memory, '
     'and whether preparations converge'),
    ('System size N, ideally two sizes', 'the finite-size reading of the ceiling s_c stands or falls with it'),
    ('Reverse loading', 'unchanged from the first pass; a parity test for a back-stress variable'),
    ('Interim model', 'the ceiling law with the one-parameter preparation memory and the aftershock hazard: '
     'the yield strength of a glass from its preparation history, to a few percent'),
]
for i, (h, t) in enumerate(steps):
    y = 1.35 + i * 1.0
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.05), Inches(0.5), Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = NAVY; sh.line.fill.background(); sh.shadow.inherit = False
    p = sh.text_frame.paragraphs[0]; run(p, str(i + 1), size=15, bold=True, color=WHITE); p.alignment = PP_ALIGN.CENTER
    sh.text_frame.margin_left = sh.text_frame.margin_right = 0; sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = textbox(s, 1.5, y, 11.2, 0.95)
    p = para(tf, True); run(p, h, size=15, bold=True, color=NAVY)
    p = tf.add_paragraph(); run(p, t, size=13, color=INK)
takeaway(s, 'Everything here is regenerated from the raw data by pipeline/run_all.sh (43 min); Figures 1 to 22 and the README carry the details.')

out = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..', 'SEEM_project_deck.pptx')
prs.save(out)
print('saved', out, len(prs.slides), 'slides')
